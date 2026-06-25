"""
Document loaders — one function per file type.

Each loader returns a list of RawPage objects, where each page/sheet/section
carries the extracted text plus source metadata that flows into the chunk payload.
"""

from __future__ import annotations

import json
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

# ─────────────────────────────────────────────
#  Data model
# ─────────────────────────────────────────────

@dataclass
class RawPage:
    """One logical 'page' extracted from a source document."""
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)


# ─────────────────────────────────────────────
#  Individual loaders
# ─────────────────────────────────────────────

def load_text(path: Path) -> list[RawPage]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return [RawPage(text=text, metadata={"page": 1})]


def load_markdown(path: Path) -> list[RawPage]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return [RawPage(text=text, metadata={"page": 1})]


def load_pdf(path: Path) -> list[RawPage]:
    try:
        import fitz  # pymupdf
    except ImportError as e:
        raise ImportError("pymupdf is required for PDF loading: pip install pymupdf") from e

    try:
        doc = fitz.open(str(path))
    except Exception:
        raise ValueError(
            f"No se pudo abrir '{path.name}'. "
            "El PDF puede estar corrompido o protegido con contraseña."
        )

    pages: list[RawPage] = []
    with doc:
        if doc.is_encrypted:
            raise ValueError(
                f"El archivo '{path.name}' está protegido con contraseña. "
                "Elimina la protección antes de ingresarlo."
            )
        for i, page in enumerate(doc, start=1):
            text = page.get_text("text")
            if text.strip():
                pages.append(RawPage(text=text, metadata={"page": i}))

    if not pages:
        raise ValueError(
            f"'{path.name}' no contiene texto extraíble. "
            "Puede ser un PDF de solo imágenes (escaneado sin OCR)."
        )
    return pages


def load_word(path: Path) -> list[RawPage]:
    try:
        from docx import Document
    except ImportError as e:
        raise ImportError("python-docx is required: pip install python-docx") from e

    try:
        doc = Document(str(path))
    except Exception:
        raise ValueError(
            f"No se pudo abrir '{path.name}'. "
            "El archivo .docx puede estar corrompido o no ser un documento Word válido."
        )
    # Preserve paragraph structure — join with newlines
    lines: list[str] = []
    current_section = 1
    pages: list[RawPage] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Detect page breaks (approximate — Word doesn't expose page numbers easily)
        lines.append(text)

    # Tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            lines.append("[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]")

    full_text = "\n".join(lines)
    return [RawPage(text=full_text, metadata={"page": 1})]


def load_excel(path: Path) -> list[RawPage]:
    try:
        import pandas as pd
    except ImportError as e:
        raise ImportError("pandas + openpyxl are required: pip install pandas openpyxl") from e

    engine = "openpyxl" if path.suffix == ".xlsx" else "xlrd"
    try:
        xl = pd.ExcelFile(str(path), engine=engine)
    except Exception:
        raise ValueError(
            f"No se pudo abrir '{path.name}'. "
            "El archivo Excel puede estar corrompido o en un formato no compatible."
        )
    pages: list[RawPage] = []

    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name)
        df = df.fillna("")
        # Convert each row to text; include header as context
        header = " | ".join(str(c) for c in df.columns)
        rows = [header]
        for _, row in df.iterrows():
            rows.append(" | ".join(str(v) for v in row.values))
        text = "\n".join(rows)
        if text.strip():
            pages.append(RawPage(
                text=text,
                metadata={"page": 1, "sheet": sheet_name},
            ))
    return pages


def load_csv(path: Path) -> list[RawPage]:
    try:
        import pandas as pd
    except ImportError as e:
        raise ImportError("pandas is required: pip install pandas") from e

    df = pd.read_csv(str(path), encoding="utf-8", errors="replace")
    df = df.fillna("")
    header = " | ".join(str(c) for c in df.columns)
    rows = [header]
    for _, row in df.iterrows():
        rows.append(" | ".join(str(v) for v in row.values))
    text = "\n".join(rows)
    return [RawPage(text=text, metadata={"page": 1})]


def load_html(path: Path) -> list[RawPage]:
    try:
        from bs4 import BeautifulSoup
    except ImportError as e:
        raise ImportError("beautifulsoup4 is required: pip install beautifulsoup4 lxml") from e

    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "lxml")
    # Remove scripts/styles
    for tag in soup(["script", "style"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return [RawPage(text=text, metadata={"page": 1})]


def load_json(path: Path) -> list[RawPage]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        raise ValueError(
            f"El archivo '{path.name}' no es un JSON válido: {e.msg} "
            f"(línea {e.lineno}, columna {e.colno})."
        )
    text = json.dumps(data, indent=2, ensure_ascii=False)
    return [RawPage(text=text, metadata={"page": 1})]


def load_pptx(path: Path) -> list[RawPage]:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError("python-pptx is required for PPTX loading: pip install python-pptx") from e

    try:
        prs = Presentation(str(path))
    except Exception:
        raise ValueError(
            f"No se pudo abrir '{path.name}'. "
            "El archivo .pptx puede estar corrompido o no ser una presentación válida."
        )

    pages: list[RawPage] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    lines.append(text)
        if lines:
            pages.append(RawPage(text="\n".join(lines), metadata={"page": i}))

    if not pages:
        raise ValueError(
            f"'{path.name}' no contiene texto extraíble. "
            "La presentación puede estar vacía o contener solo imágenes."
        )
    return pages


# ─────────────────────────────────────────────
#  Dispatcher
# ─────────────────────────────────────────────

_LOADERS = {
    "text":     load_text,
    "markdown": load_markdown,
    "pdf":      load_pdf,
    "word":     load_word,
    "excel":    load_excel,
    "csv":      load_csv,
    "html":     load_html,
    "json":     load_json,
    "pptx":     load_pptx,
}


def load_document(path: Path, loader_key: str) -> list[RawPage]:
    """
    Load a document from *path* using the appropriate loader.

    Returns a list of RawPage objects (one per PDF page / Excel sheet / etc.).
    """
    loader = _LOADERS.get(loader_key)
    if loader is None:
        raise ValueError(f"No loader registered for key '{loader_key}'")
    pages = loader(path)
    # Attach universal source metadata to every page
    for page in pages:
        page.metadata.update({
            "source_file": path.name,
            "source_path": str(path),
            "file_type":   path.suffix.lstrip(".").lower(),
        })
    return pages
