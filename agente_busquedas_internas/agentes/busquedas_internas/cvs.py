"""
Servicio de lectura y búsqueda semántica de CVs en formato PPTX.

el flujo es:
  1. Leer todos los .pptx del directorio de CVs (por default: datos_prueba/cvs/).
  2. Generar embeddings con sentence-transformers (modelo local, sin API key).
  3. Cachea embeddings en disco (.embeddings_cache.json junto a los CVs).
     El caché se invalida por archivo cuando su mtime cambia, por lo que
     un CV modificado genera un nuevo embedding automáticamente.
  4. Dado un texto de JD, devuelve los N CVs más similares por similitud coseno.

Este módulo es la fuente de datos local mientras la integración con Workday
no esté disponible. Cuando Workday esté listo, solo consultar_ats.py cambia;
este módulo puede coexistir como fuente alternativa o de prueba.

Dependencias externas: python-pptx, numpy, sentence-transformers.
El modelo se descarga automáticamente la primera vez (~90MB, sin API key).
"""
import json
from pathlib import Path

import numpy as np
from dotenv import load_dotenv
from pptx import Presentation
from sentence_transformers import SentenceTransformer

from agentes.config.settings import get_settings
from schemas.cvs_data import Cvs_data

load_dotenv()

# region Configuración de rutas

# CVS_DIR y CACHE_FILE se resuelven la primera vez que se llama a Buscar_candidatos_similares,
# no al importar el módulo. Esto evita que un .env ausente rompa el import y bloquee
# los breakpoints del debugger antes de entrar a cualquier función.
_st_model: SentenceTransformer | None = None


def _get_model() -> SentenceTransformer:
    global _st_model
    if _st_model is None:
        _st_model = SentenceTransformer(get_settings().EMBEDDING_MODEL)
    return _st_model
# endregion

# region Extracción de texto

def Extraer_texto_pptx(path: Path) -> str:
    """
    Extrae y concatena todo el texto de un archivo .pptx.

    Los CVs de Accenture tienen 2 slides (inglés + español del mismo contenido).
    Se extrae texto de todas las slides. La redundancia bilingüe no perjudica
    la calidad del embedding: al contrario, refuerza los conceptos clave
    al aparecer en dos idiomas.
    """
    prs = Presentation(str(path))
    fragmentos = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                fragmentos.append(shape.text.strip())
    return "\n".join(fragmentos)
# endregion

# region Cache de embeddings

def Cargar_cache(cache_file: Path) -> dict:
    """Carga el caché desde disco. Retorna dict vacío si no existe todavía."""
    if cache_file.exists():
        with open(cache_file, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def _guardar_cache(cache: dict, cache_file: Path) -> None:
    """Persiste el caché a disco (sobreescribe el archivo existente)."""
    with open(cache_file, "w", encoding="utf-8") as f:
        json.dump(cache, f, ensure_ascii=False)

# endregion

# region Embeddings y similitud
def Generar_embedding(texto: str) -> list[float]:
    """Genera el vector de embedding de un texto usando el modelo local."""
    return _get_model().encode(texto).tolist()


def Batch_cosine_similarity(
    query: list[float],
    vectores: list[list[float]],
) -> list[float]:
    """
    Calcula la similitud coseno del vector query contra todos los vectores
    en una sola operación matricial (numpy), sin iterar uno a uno.

    Esto escala eficientemente a miles de CVs: una matriz de 10.000 × 768
    tarda ~5ms en numpy vs ~500ms con un loop Python puro.

    Retorna una lista de floats (0.0–1.0) en el mismo orden que `vectores`.
    """
    q = np.array(query, dtype=np.float32)
    matrix = np.array(vectores, dtype=np.float32)

    # Normalizar query y filas de la matriz a norma unitaria
    q_norm = q / (np.linalg.norm(q) + 1e-10)
    norms = np.linalg.norm(matrix, axis=1, keepdims=True) + 1e-10
    matrix_norm = matrix / norms

    # Producto punto vectorizado = similitud coseno para vectores normalizados
    return matrix_norm.dot(q_norm).tolist()

# endregion


def Buscar_candidatos_similares(
    jd_texto: str,
    top_n: int = 10,
) -> list[Cvs_data]:
    """
    Dado un texto de Job Description, devuelve los N CVs más similares
    usando búsqueda semántica por embeddings.

    Los embeddings de los CVs se cachean en disco. Si un .pptx no cambió
    desde la última ejecución (mismo mtime), se reutiliza el embedding cacheado
    sin llamar a la API. Solo los CVs nuevos o modificados generan una llamada.

    Args:
        jd_texto: Texto de la JD (role_title + skills + management_level).
        top_n: Número máximo de CVs a retornar. Default: 10.

    Returns:
        Lista de dicts ordenada de mayor a menor similitud semántica:
            [
                {
                    "id": str,               # nombre del archivo sin extensión (ej: "Julian Monte")
                    "nombre": str,           # ídem, para legibilidad del LLM
                    "texto_cv": str,         # texto completo extraído del .pptx
                    "score_embedding": float # similitud coseno con la JD (0.0–1.0)
                },
                ...
            ]
    """
    cvs_dir = get_settings().CVS_DIR
    cache_file = cvs_dir / ".embeddings_cache.json"

    cache = Cargar_cache(cache_file)
    cache_actualizado = False
    cvs_data: list[Cvs_data] = []

    # --- Paso 1: cargar / actualizar embeddings de los CVs ---
    for pptx_path in sorted(cvs_dir.glob("*.pptx")):
        nombre = pptx_path.stem
        mtime = str(pptx_path.stat().st_mtime) # timestamp de última modificación

        entrada = cache.get(nombre)
        if entrada and entrada.get("mtime") == mtime:
            # CV sin cambios: reusar embedding cacheado (sin llamada a la API)
            embedding = entrada["embedding"]
            texto = entrada["texto"]
        else:
            # CV nuevo o modificado: re-extraer texto y generar embedding
            texto = Extraer_texto_pptx(pptx_path)
            embedding = Generar_embedding(texto)
            cache[nombre] = {"mtime": mtime, "embedding": embedding, "texto": texto}
            cache_actualizado = True

        cvs_data.append(
            Cvs_data(
            id=nombre,
            nombre=nombre,
            texto_cv=texto,
            embedding=embedding
        ))

    if cache_actualizado:
        _guardar_cache(cache, cache_file)

    if not cvs_data:
        return []

    # --- Paso 2: embedding de la JD ---
    jd_embedding : list[float] = Generar_embedding(jd_texto)

    # --- Paso 3: similitud coseno vectorizada y ordenamiento ---
    vectores = [cv.embedding for cv in cvs_data]
    scores = Batch_cosine_similarity(jd_embedding, vectores)

    # --- calculo el score para cada candidato y lo guardo en el objeto ---
    for cv, score in zip(cvs_data, scores):
        cv.score_embedding = round(float(score), 4)

    cvs_data.sort(key=lambda x: x.score_embedding, reverse=True)

    # Retornar top N sin el vector de embedding (innecesario para el LLM downstream)
    return [
        Cvs_data(
            id=cv.nombre,
            nombre=cv.nombre,
            texto_cv=cv.texto_cv,
            score_embedding=cv.score_embedding
        )
        for cv in cvs_data[:top_n]
    ]
